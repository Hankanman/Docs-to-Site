"""
Core document conversion functionality.
"""
import logging
import os
import re
import shutil
from pathlib import Path
from typing import List, Optional, Set, Tuple, Dict, Any
from unicodedata import normalize
from slugify import slugify

import yaml
from markitdown import MarkItDown, UnsupportedFormatException, FileConversionException
from pptx import Presentation
from pptx.shapes.picture import Picture

logger = logging.getLogger(__name__)

# Supported input formats (all lowercase)
SUPPORTED_FORMATS = {
    # Documents
    '.pdf', '.doc', '.docx', '.rtf', '.odt', '.txt',
    # Presentations
    '.ppt', '.pptx',
    # Spreadsheets
    '.xls', '.xlsx', '.csv',
    # Images
    '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp',
    # Audio
    '.mp3', '.wav', '.m4a', '.ogg', '.flac',
    # Web
    '.html', '.htm',
    # Data formats
    '.json', '.xml',
    # Archives
    '.zip'
}

# Image formats that should be copied to the output
IMAGE_FORMATS = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.svg'}


def sanitize_filename(filename: str) -> str:
    """
    Sanitize a filename to be safe for all operating systems.
    
    Args:
        filename: The filename to sanitize
    
    Returns:
        A sanitized filename
    """
    # Get the stem and extension separately
    path = Path(filename)
    # Handle paths with multiple segments
    if len(path.parts) > 1:
        # Join all parts except the last with hyphens
        stem = '-'.join(part for part in path.parts[:-1])
        stem += f"-{path.parts[-1]}"
    else:
        stem = path.stem
    suffix = path.suffix
    
    # Slugify the stem only
    clean_stem = slugify(stem,
                        lowercase=False,
                        separator='-',
                        regex_pattern=r'[^-a-zA-Z0-9.]+',  # Allow dots
                        replacements=[
                            ('_', '-'),    # Convert underscores to hyphens
                            ('/', '-'),    # Convert slashes to hyphens
                            ('\\', '-'),   # Convert backslashes to hyphens
                        ])
    
    # Recombine with the extension
    return f"{clean_stem}{suffix}"


def sanitize_title(title: str) -> str:
    """
    Sanitize a title for use in MkDocs navigation.
    
    Args:
        title: The title to sanitize
    
    Returns:
        A sanitized title
    """
    # First handle special characters and spacing
    replacements = {
        '–': '-',  # en dash
        '—': '-',  # em dash
        '™': '',   # trademark
        '®': '',   # registered trademark
        '©': '',   # copyright
        '[': '(',  # brackets to parentheses
        ']': ')',
    }
    
    # Apply replacements
    clean_title = title
    for old, new in replacements.items():
        clean_title = clean_title.replace(old, new)
    
    # Remove other special characters but keep unicode letters/numbers
    clean_title = ''.join(c for c in clean_title if c.isalnum() or c.isspace() or c in '()-.,')
    
    # Normalize spaces
    clean_title = ' '.join(clean_title.split())
    
    return clean_title.strip()


class DocumentConverter:
    """Handles the conversion of documents to Markdown and MkDocs site generation."""

    def __init__(self, input_dir: Path, output_dir: Path, config: Optional[Path] = None) -> None:
        """
        Initialize the document converter.

        Args:
            input_dir: Directory containing input documents
            output_dir: Directory where the MkDocs site will be generated
            config: Optional path to custom MkDocs configuration
        """
        self.input_dir = input_dir
        self.output_dir = output_dir
        self.config = config
        self.docs_dir = output_dir / 'docs'
        self.mkdocs_config = output_dir / 'mkdocs.yml'
        self.converter = MarkItDown()
        self.converted_files: Dict[Path, str] = {}  # Maps output paths to titles
        self.images_dir = self.docs_dir / 'images'

    def is_supported_format(self, file_path: Path) -> bool:
        """Check if a file has a supported format."""
        return file_path.suffix.lower() in SUPPORTED_FORMATS

    def is_image_format(self, file_path: Path) -> bool:
        """Check if a file is an image format."""
        return file_path.suffix.lower() in IMAGE_FORMATS

    def get_documents(self) -> List[Tuple[Path, bool]]:
        """
        Find all supported documents in the input directory.
        
        Returns:
            List of tuples containing (file_path, is_accessible)
        """
        documents = []
        for file in self.input_dir.rglob('*'):
            if file.is_file() and self.is_supported_format(file):
                # Check if file is accessible
                try:
                    with open(file, 'rb') as f:
                        f.read(1)
                    documents.append((file, True))
                    logger.info(f"Found supported document: {file}")
                except (PermissionError, OSError):
                    logger.warning(f"File {file} exists but cannot be accessed. It may be open in another program.")
                    documents.append((file, False))

        logger.info(f"Found {len(documents)} supported documents")
        if not documents:
            logger.warning(f"No supported documents found. Supported formats are: {', '.join(sorted(SUPPORTED_FORMATS))}")
        return documents

    def normalize_image_name(self, name: str) -> str:
        """Normalize image name to handle variations in spacing and extensions."""
        # Remove extension
        base = str(Path(name).stem)
        # Remove spaces
        base = re.sub(r'\s+', '', base)
        # Convert to lowercase for case-insensitive matching
        return base.lower()

    def extract_pptx_images(self, document: Path, doc_images_dir: Path) -> Dict[str, str]:
        """
        Extract images from a PowerPoint file.
        
        Args:
            document: Path to the PowerPoint file
            doc_images_dir: Directory to save images in
        
        Returns:
            Dictionary mapping original image filenames to new image paths
        """
        image_map = {}
        try:
            prs = Presentation(document)
            image_count = 0

            for slide in prs.slides:
                for shape in slide.shapes:
                    if isinstance(shape, Picture):
                        try:
                            # Get image data and format
                            image = shape.image
                            image_bytes = image.blob
                            image_ext = image.ext.lower()
                            if not image_ext.startswith('.'):
                                image_ext = '.' + image_ext

                            # Generate unique filename
                            image_count += 1
                            image_name = f"image_{image_count}{image_ext}"
                            image_path = doc_images_dir / image_name

                            # Save image
                            with open(image_path, 'wb') as f:
                                f.write(image_bytes)
                            logger.info(f"Extracted image: {image_path}")

                            # Map all possible variations of the image name
                            possible_names = []
                            if hasattr(shape, 'name'):
                                possible_names.append(shape.name)
                            if hasattr(shape, 'image_filename'):
                                possible_names.append(shape.image_filename)
                            
                            # Add common PowerPoint image naming patterns
                            for name in possible_names:
                                # Store the full name
                                image_map[name] = image_name
                                # Store just the filename
                                image_map[Path(name).name] = image_name
                                # Store normalized name for matching
                                norm_name = self.normalize_image_name(name)
                                image_map[norm_name] = image_name
                                # Add common extensions
                                for ext in ['.jpg', '.jpeg', '.png', '.gif', '.wmf']:
                                    image_map[f"{norm_name}{ext}"] = image_name
                                
                            logger.debug(f"Image mappings for {image_name}: {[k for k in image_map.keys() if image_map[k] == image_name]}")

                        except Exception as e:
                            logger.warning(f"Failed to extract image from shape in {document}: {str(e)}")
                            logger.debug(f"Shape details - name: {getattr(shape, 'name', 'N/A')}, type: {type(shape)}")

        except Exception as e:
            logger.warning(f"Failed to extract images from {document}: {str(e)}")

        logger.info(f"Total images extracted: {image_count}")
        logger.info(f"Image mappings: {image_map}")
        return image_map

    def format_markdown(self, content: str, document: Path, image_map: Optional[Dict[str, str]] = None) -> str:
        """
        Format the Markdown content for better readability.
        
        Args:
            content: Original Markdown content
            document: Path to the source document
            image_map: Optional mapping of original image names to new paths
        
        Returns:
            Formatted Markdown content
        """
        # Format slide markers
        content = re.sub(r'<!-- Slide number: (\d+) -->', r'\n---\n### Slide \1\n', content)
        
        # Update image paths if we have a mapping
        if image_map:
            # Log all image references in the markdown
            image_refs = re.findall(r'!\[(.*?)\]\((.*?)\)', content)
            logger.info(f"Found image references in markdown: {image_refs}")
            
            # First try to match exact filenames
            for old_name, new_path in image_map.items():
                pattern = rf'!\[(.*?)\]\({re.escape(old_name)}\)'
                if re.search(pattern, content):
                    logger.debug(f"Found match for {old_name} -> {new_path}")
                content = re.sub(
                    pattern,
                    rf'![\1](images/{sanitize_filename(document.stem)}/{new_path})',
                    content
                )

            # Then try to match using normalized names
            for match in re.finditer(r'!\[(.*?)\]\((.*?)\)', content):
                alt_text = match.group(1)
                img_name = match.group(2)
                norm_name = self.normalize_image_name(img_name)
                
                # Try to find a match using the normalized name
                if norm_name in image_map:
                    new_path = image_map[norm_name]
                    logger.debug(f"Found normalized match: {norm_name} -> {new_path}")
                    content = content.replace(
                        match.group(0),
                        f'![{alt_text}](images/{sanitize_filename(document.stem)}/{new_path})'
                    )
                else:
                    logger.warning(f"No mapping found for image: {img_name} (normalized: {norm_name})")
                    # Don't modify the path if we can't find a mapping
                    continue

        # Fix any remaining absolute image paths
        content = re.sub(
            r'!\[(.*?)\]\(/images/',
            r'![\1](images/',
            content
        )
        
        # Remove any duplicate image paths
        content = re.sub(
            r'images/[^/]+/images/[^/]+/',
            lambda m: m.group(0).split('/', 1)[1],
            content
        )
        
        # Log any remaining unmatched image references
        remaining_refs = re.findall(r'!\[(.*?)\]\((.*?)\)', content)
        logger.info(f"Remaining image references after processing: {remaining_refs}")
        
        # Remove extra blank lines
        content = re.sub(r'\n{3,}', r'\n\n', content)
        
        return content

    def copy_embedded_images(self, document: Path, doc_images_dir: Path) -> Dict[str, str]:
        """
        Copy embedded images from the document's directory to the output images directory.
        
        Args:
            document: Path to the document being converted
            doc_images_dir: Directory to copy images to
            
        Returns:
            Dictionary mapping original image filenames to new paths
        """
        image_map = {}
        # Look for images in the same directory as the document
        doc_dir = document.parent
        for img_ext in IMAGE_FORMATS:
            for img_file in doc_dir.glob(f"*{img_ext}"):
                try:
                    # Generate a sanitized name for the image
                    sanitized_name = sanitize_filename(img_file.name)
                    new_path = doc_images_dir / sanitized_name
                    
                    # Copy the image
                    shutil.copy2(img_file, new_path)
                    logger.info(f"Copied embedded image: {img_file} -> {new_path}")
                    
                    # Store both the original name and sanitized name in the map
                    image_map[img_file.name] = sanitized_name
                    image_map[sanitized_name] = sanitized_name
                except Exception as e:
                    logger.warning(f"Failed to copy embedded image {img_file}: {str(e)}")
        
        return image_map

    def convert_document(self, document: Path) -> Path:
        """
        Convert a single document to Markdown.

        Args:
            document: Path to the document to convert

        Returns:
            Path to the converted Markdown file

        Raises:
            FileConversionException: If the file cannot be converted
            PermissionError: If the file cannot be accessed
            OSError: If there are file system related errors
        """
        relative_path = document.relative_to(self.input_dir)
        # Sanitize the filename part while keeping the directory structure
        sanitized_name = sanitize_filename(relative_path.stem) + '.md'
        sanitized_path = relative_path.parent / sanitized_name
        output_path = self.docs_dir / sanitized_path
        
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        logger.info(f"Converting {document} to {output_path}")
        try:
            # First check if we can access the file
            with open(document, 'rb') as f:
                f.read(1)
            
            # Create document-specific images directory in the same directory as the markdown file
            doc_images_dir = output_path.parent / 'images' / sanitize_filename(document.stem)
            doc_images_dir.mkdir(parents=True, exist_ok=True)
            
            # Initialize image map
            image_map = {}
            
            # Copy any embedded images from the document's directory
            embedded_images = self.copy_embedded_images(document, doc_images_dir)
            if embedded_images:
                image_map.update(embedded_images)
            
            # Extract images if it's a PowerPoint file
            if document.suffix.lower() in {'.ppt', '.pptx'}:
                pptx_images = self.extract_pptx_images(document, doc_images_dir)
                if pptx_images:
                    image_map.update(pptx_images)
            
            # Convert document to Markdown
            result = self.converter.convert_local(str(document))
            
            # Get the title and content
            title = getattr(result, 'title', None)
            if not title:
                # Use the filename without extension as title if no title is found
                title = document.stem
            title = sanitize_title(title)
            
            content = result.text_content
            
            # If we have a title, add it as a header
            markdown_content = f"# {title}\n\n" if title else ""
            markdown_content += content
            
            # Format the content
            markdown_content = self.format_markdown(markdown_content, document, image_map)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            # Store the title for navigation
            relative_output = output_path.relative_to(self.docs_dir)
            self.converted_files[relative_output] = title
            
            return output_path
        except (PermissionError, OSError) as e:
            logger.error(f"Cannot access file {document}. The file may be open in another program or you may not have permission to access it.")
            raise
        except UnsupportedFormatException as e:
            logger.error(f"Unsupported format for file: {document}")
            raise
        except FileConversionException as e:
            logger.error(f"Error converting {document}: {str(e)}")
            raise
        except Exception as e:
            logger.error(f"Unexpected error converting {document}: {str(e)}")
            raise

    def _build_nav_structure(self) -> List[Any]:
        """Build the navigation structure for MkDocs."""
        nav_structure: Dict[str, Any] = {}

        # Sort files by path for consistent ordering
        sorted_files = sorted(self.converted_files.items(), key=lambda x: str(x[0]))

        # Group files by prefix (e.g., "Client -" or "watsonx")
        groups: Dict[str, List[Tuple[Path, str]]] = {}
        for file_path, title in sorted_files:
            # Extract prefix from title
            prefix = title.split(' - ')[0] if ' - ' in title else 'Other'
            # Sanitize the prefix
            prefix = sanitize_title(prefix)
            if prefix not in groups:
                groups[prefix] = []
            groups[prefix].append((file_path, title))

        # Build navigation structure with groups
        for prefix, files in groups.items():
            if len(files) > 1:
                # Create a group for multiple files with the same prefix
                nav_structure[prefix] = {
                    sanitize_title(title.replace(prefix + ' - ', '')): str(file_path)
                    for file_path, title in files
                }
            else:
                # Single file goes directly in the root
                file_path, title = files[0]
                nav_structure[sanitize_title(title)] = str(file_path)

        # Convert the nested dictionary to MkDocs nav format
        def dict_to_nav(d: Dict[str, Any]) -> List[Any]:
            nav = []
            for k, v in d.items():
                if isinstance(v, str):
                    nav.append({k: v})
                else:
                    # Handle groups
                    subnav = []
                    for sub_k, sub_v in v.items():
                        subnav.append({sub_k: sub_v})
                    nav.append({k: subnav})
            return nav

        return dict_to_nav(nav_structure)

    def generate_mkdocs_config(self) -> None:
        """Generate the MkDocs configuration file with navigation structure."""
        # Ensure output directory exists
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        if self.config:
            # Use custom config if provided
            with open(self.config, 'r', encoding='utf-8') as f:
                config_data = yaml.safe_load(f)
        else:
            # Generate default config
            config_data = {
                'site_name': 'Documentation',
                'theme': {
                    'name': 'material',
                    'features': [
                        'navigation.instant',
                        'navigation.tracking',
                        'navigation.sections',
                        'navigation.expand',
                        'toc.integrate'
                    ],
                    'palette': {
                        'primary': 'blue',
                        'accent': 'blue'
                    }
                },
                'docs_dir': 'docs',
                'use_directory_urls': False,  # This helps with relative image paths
                'markdown_extensions': [
                    'attr_list',
                    'md_in_html',
                    'tables',
                    'fenced_code',
                    'footnotes',
                    'admonition',
                    'toc',
                    'pymdownx.highlight',
                    'pymdownx.superfences'
                ]
            }

        # Add navigation structure
        if self.converted_files:
            config_data['nav'] = self._build_nav_structure()

        with open(self.mkdocs_config, 'w', encoding='utf-8') as f:
            yaml.dump(
                config_data, 
                f, 
                default_flow_style=False, 
                sort_keys=False, 
                allow_unicode=True,
                encoding='utf-8'
            )

    def convert(self) -> None:
        """Convert all documents and set up the MkDocs site."""
        logger.info(f"Starting document conversion from {self.input_dir} to {self.output_dir}")

        # Create output directory
        self.docs_dir.mkdir(parents=True, exist_ok=True)

        # Convert documents
        documents = self.get_documents()
        if not documents:
            raise ValueError(f"No supported documents found in {self.input_dir}")

        conversion_errors = []
        inaccessible_files = []

        for doc, is_accessible in documents:
            if not is_accessible:
                inaccessible_files.append(doc)
                continue

            try:
                self.convert_document(doc)
            except (PermissionError, OSError):
                inaccessible_files.append(doc)
            except Exception as e:
                logger.error(f"Failed to convert {doc}: {str(e)}")
                conversion_errors.append((doc, str(e)))

        # Generate MkDocs configuration
        self.generate_mkdocs_config()

        # Report conversion summary
        total = len(documents)
        inaccessible = len(inaccessible_files)
        failed = len(conversion_errors)
        succeeded = total - failed - inaccessible

        logger.info(f"Document conversion complete:")
        logger.info(f"- Successfully converted: {succeeded}/{total} documents")
        if inaccessible_files:
            logger.warning("The following files could not be accessed (they may be open in another program):")
            for doc in inaccessible_files:
                logger.warning(f"- {doc}")
        if conversion_errors:
            logger.warning("The following documents failed to convert:")
            for doc, error in conversion_errors:
                logger.warning(f"- {doc}: {error}") 