"""
Image processing functionality for document conversion.
"""

import logging
import shutil
import subprocess
from pathlib import Path
from typing import Dict, Optional, Tuple

from pptx import Presentation
from pptx.shapes.picture import Picture

from .utils import (
    sanitize_filename,
    IMAGE_FORMATS,
    normalize_image_name,
    check_imagemagick,
)

logger = logging.getLogger(__name__)


def convert_wmf_to_png(
    image_bytes: bytes, output_path: Path, is_icon: bool = True
) -> Optional[Path]:
    """
    Convert WMF image to PNG or SVG format using ImageMagick.
    Attempts to preserve vector content by using direct ImageMagick commands.

    Args:
        image_bytes: The WMF image data
        output_path: The desired output path for the converted file
        is_icon: Whether the WMF is being used as an icon (affects size/density settings)

    Returns:
        Path to the converted file or None if conversion fails
    """
    # Check ImageMagick availability first
    is_available, message = check_imagemagick()
    if not is_available:
        logger.warning(f"Cannot convert WMF: {message}")
        return None

    try:
        # Create a temporary WMF file
        temp_wmf = output_path.with_suffix(".tmp.wmf")
        temp_wmf.write_bytes(image_bytes)

        try:
            # Try vector conversion first using direct ImageMagick command
            svg_path = output_path.with_suffix(".svg")

            # Adjust settings based on whether it's an icon
            density = "72" if is_icon else "300"
            size_param = (
                ["-resize", "64x64>"] if is_icon else []
            )  # Resize if larger than 64x64

            convert_cmd = [
                "magick",
                "convert",
                "-density",
                density,
                "-background",
                "transparent",
                "-define",
                "wmf:vector",
                "-define",
                "wmf:preserve-vector",
                "-define",
                "svg:include-xml",
                str(temp_wmf),
            ]

            # Add size parameter for icons
            if size_param:
                convert_cmd.extend(size_param)

            convert_cmd.append("SVG:" + str(svg_path))

            result = subprocess.run(convert_cmd, capture_output=True, text=True)
            if result.returncode == 0:
                # Verify the SVG is not just a base64 wrapper
                svg_content = svg_path.read_text()
                if "base64" not in svg_content and (
                    "<path" in svg_content
                    or "<polygon" in svg_content
                    or "<rect" in svg_content
                ):
                    logger.info(f"Successfully converted WMF to SVG: {svg_path}")
                    return svg_path
                else:
                    logger.warning(
                        "SVG conversion resulted in base64 encoding, falling back to PNG"
                    )
                    svg_path.unlink(missing_ok=True)

            # Fallback to PNG conversion
            png_path = output_path.with_suffix(".png")
            convert_cmd = [
                "magick",
                "convert",
                "-density",
                density,
                "-background",
                "white",
                "-alpha",
                "remove",
                str(temp_wmf),
            ]

            # Add size parameter for icons
            if size_param:
                convert_cmd.extend(size_param)

            convert_cmd.append(str(png_path))

            result = subprocess.run(convert_cmd, capture_output=True, text=True)
            if result.returncode == 0:
                logger.info(f"Successfully converted WMF to PNG: {png_path}")
                return png_path
            else:
                logger.warning(f"PNG conversion failed: {result.stderr}")
                return None

        finally:
            # Clean up temporary file
            temp_wmf.unlink(missing_ok=True)

    except Exception as e:
        logger.warning(f"Failed to convert WMF: {str(e)}")
        return None


def process_image_shape(shape: Picture, image_count: int) -> Tuple[str, bytes, str]:
    """
    Process a PowerPoint Picture shape and extract its image data.

    Args:
        shape: PowerPoint Picture shape
        image_count: Current image counter for naming

    Returns:
        Tuple of (image_name, image_bytes, image_ext)
    """
    image = shape.image
    image_bytes = image.blob
    image_ext = image.ext.lower()

    # For GIF files, keep the original extension
    if image_ext == "gif":
        image_name = f"image_{image_count}.gif"
    # For WMF files, always convert to PNG
    elif image_ext == "wmf":
        image_name = f"image_{image_count}.png"
    else:
        image_name = f"image_{image_count}.{image_ext}"

    logger.debug(f"Processing image shape: {image_name} (original name: {getattr(shape, 'name', 'N/A')})")
    return image_name, image_bytes, image_ext


def create_image_mappings(shape: Picture, image_name: str) -> Dict[str, str]:
    """
    Create all possible mappings for an image name.

    Args:
        shape: PowerPoint Picture shape
        image_name: The target image name

    Returns:
        Dictionary of image name mappings
    """
    mappings = {}
    possible_names = []

    if hasattr(shape, "name"):
        possible_names.append(shape.name)
    if hasattr(shape, "image_filename"):
        possible_names.append(shape.image_filename)

    # Add mappings for all possible names
    for name in possible_names:
        # Store the full name
        mappings[name] = image_name
        mappings[Path(name).name] = image_name
        mappings[normalize_image_name(name)] = image_name
        # Add common extensions (no WMF)
        for ext in [".jpg", ".jpeg", ".png", ".gif"]:
            mappings[f"{normalize_image_name(name)}{ext}"] = image_name
            mappings[f"{Path(name).stem}{ext}"] = image_name

    return mappings


def find_unique_images(
    prs: Presentation,
) -> Tuple[Dict[int, Tuple[int, int]], Dict[int, str], int]:
    """
    First pass through presentation to find unique images and their first occurrences.

    Args:
        prs: PowerPoint presentation object

    Returns:
        Tuple of (first_occurrence dict, blob_to_image dict, image_count)
    """
    first_occurrence = {}  # Maps blob hash to slide/shape where first seen
    blob_to_image = {}  # Maps blob hash to first image name for that blob
    image_count = 0  # Start at 0, will increment before assigning names
    processed_hashes = set()  # Track which hashes we've seen to maintain order
    ordered_hashes = []  # Track the order of first appearance

    # First pass: Find all unique images and their first occurrences
    for slide_num, slide in enumerate(prs.slides, 1):
        logger.debug(f"Processing slide {slide_num} for unique images")
        for shape_idx, shape in enumerate(slide.shapes, 1):
            if isinstance(shape, Picture):
                try:
                    image = shape.image
                    image_bytes = image.blob
                    blob_hash = hash(image_bytes)

                    # Track the order of first appearance for each unique image
                    if blob_hash not in processed_hashes:
                        processed_hashes.add(blob_hash)
                        ordered_hashes.append(blob_hash)
                        first_occurrence[blob_hash] = (slide_num, shape_idx)
                        logger.debug(
                            f"Found new unique image in slide {slide_num}, "
                            f"shape {shape_idx} (hash: {blob_hash})"
                        )
                except Exception as e:
                    logger.warning(f"Failed to process image in first pass: {str(e)}")

    # Second pass: Assign image names in order of appearance
    for image_count, blob_hash in enumerate(ordered_hashes, 1):
        slide_num, shape_idx = first_occurrence[blob_hash]
        try:
            # Get the image extension from the original shape
            shape = prs.slides[slide_num-1].shapes[shape_idx-1]
            ext = shape.image.ext.lower()
            # Handle WMF files differently - they will be converted to PNG
            if ext == "wmf":
                image_name = f"image_{image_count}.png"
            else:
                image_name = f"image_{image_count}.{ext}"
            blob_to_image[blob_hash] = image_name
            logger.debug(
                f"Assigned name {image_name} to image from slide {slide_num}, "
                f"shape {shape_idx} (hash: {blob_hash})"
            )
        except Exception as e:
            logger.warning(f"Failed to assign image name: {str(e)}")

    logger.info(f"Found {len(ordered_hashes)} unique images across {len(prs.slides)} slides")
    logger.debug(f"First occurrence mapping: {first_occurrence}")
    logger.debug(f"Blob to image mapping: {blob_to_image}")
    return first_occurrence, blob_to_image, len(ordered_hashes)


def extract_pptx_images(document: Path, doc_images_dir: Path) -> Dict[str, str]:
    """
    Extract images from a PowerPoint file.

    Args:
        document: Path to the PowerPoint file
        doc_images_dir: Directory to save images in

    Returns:
        Dictionary mapping original image filenames to new image paths
    """
    # Check ImageMagick at the start
    is_available, message = check_imagemagick()
    if not is_available:
        logger.warning(
            f"ImageMagick is not available: {message}. WMF images will not be converted to PNG."
        )

    image_map = {}
    processed_images = set()  # Track which images have been processed

    try:
        # First pass: Find all unique images and their first occurrences
        prs = Presentation(document)
        first_occurrence, blob_to_image, _ = find_unique_images(prs)

        logger.debug(f"First occurrence mapping: {first_occurrence}")
        logger.debug(f"Blob to image mapping: {blob_to_image}")

        # Second pass: Extract images and build mappings in order of appearance
        for slide_num, slide in enumerate(prs.slides, 1):
            logger.debug(f"Processing slide {slide_num} for image extraction")
            for shape_idx, shape in enumerate(slide.shapes, 1):
                if isinstance(shape, Picture):
                    try:
                        # Get image data
                        image = shape.image
                        image_bytes = image.blob
                        blob_hash = hash(image_bytes)
                        image_name = blob_to_image[blob_hash]
                        image_ext = image.ext.lower()

                        logger.debug(
                            f"Processing image in slide {slide_num}, shape {shape_idx}: "
                            f"shape_name={getattr(shape, 'name', 'N/A')}, "
                            f"image_name={image_name}, hash={blob_hash}"
                        )

                        # Only extract the image if we haven't processed it yet
                        if blob_hash not in processed_images:
                            processed_images.add(blob_hash)
                            image_path = doc_images_dir / image_name

                            # Handle WMF files - convert to PNG if ImageMagick is available
                            if image_ext == "wmf":
                                if is_available:
                                    logger.debug(f"Converting WMF image to PNG: {image_name}")
                                    png_path = convert_wmf_to_png(
                                        image_bytes, image_path, is_icon=True
                                    )
                                    if not png_path:
                                        logger.warning(
                                            f"Failed to convert WMF to PNG: {image_name}. Saving original WMF."
                                        )
                                        # Save original WMF if conversion fails
                                        wmf_path = image_path.with_suffix(".wmf")
                                        with open(wmf_path, "wb") as f:
                                            f.write(image_bytes)
                                else:
                                    logger.warning(
                                        f"ImageMagick not available. Saving original WMF: {image_name}"
                                    )
                                    # Save as WMF since we can't convert
                                    wmf_path = image_path.with_suffix(".wmf")
                                    with open(wmf_path, "wb") as f:
                                        f.write(image_bytes)
                            else:
                                # Save other image types
                                logger.debug(f"Saving image: {image_path}")
                                with open(image_path, "wb") as f:
                                    f.write(image_bytes)

                            logger.info(f"Extracted image: {image_path}")
                        else:
                            logger.debug(
                                f"Skipping duplicate image extraction for {image_name}"
                            )

                        # Update image mappings for this shape
                        shape_mappings = create_image_mappings(shape, image_name)
                        logger.debug(f"Created mappings for {image_name}: {shape_mappings}")
                        image_map.update(shape_mappings)

                    except Exception as e:
                        logger.warning(
                            f"Failed to extract image from shape in {document}: {str(e)}"
                        )
                        logger.debug(
                            f"Shape details - name: {getattr(shape, 'name', 'N/A')}, type: {type(shape)}"
                        )

    except Exception as e:
        logger.warning(f"Failed to extract images from {document}: {str(e)}")

    logger.info(f"Total unique images extracted: {len(processed_images)}")
    logger.info(f"Total image mappings: {len(image_map)}")
    logger.debug(f"Final image mappings: {image_map}")
    return image_map


def copy_embedded_images(document: Path, doc_images_dir: Path) -> Dict[str, str]:
    """
    Copy embedded images from the document's directory to the output images directory.

    Args:
        document: Path to the document being converted
        doc_images_dir: Directory to copy images to

    Returns:
        Dictionary mapping original image filenames to new paths
    """
    image_map = {}
    # Look for images in the same directory as the document
    doc_dir = document.parent
    for img_ext in IMAGE_FORMATS:
        for img_file in doc_dir.glob(f"*{img_ext}"):
            try:
                # Generate a sanitized name for the image
                sanitized_name = sanitize_filename(img_file.name)
                new_path = doc_images_dir / sanitized_name

                # Copy the image
                shutil.copy2(img_file, new_path)
                logger.info(f"Copied embedded image: {img_file} -> {new_path}")

                # Store both the original name and sanitized name in the map
                image_map[img_file.name] = sanitized_name
                image_map[sanitized_name] = sanitized_name
            except Exception as e:
                logger.warning(f"Failed to copy embedded image {img_file}: {str(e)}")

    return image_map
